import streamlit as st
import os
import json
import re
import qrcode
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# App-Konfiguration & Design
st.set_page_config(
    page_title="Evasys Schritt 9 - Universal Unterlagen Ersteller",
    page_icon="🚀",
    layout="wide"
)

st.title("🚀 Schritt 9: Präsentationsunterlagen & Ordnerstrukturen generieren")
st.markdown("""
Dieses Tool liest die hochgeladenen Strukturdaten, generiert QR-Codes, erstellt die HTML-Dateien (RLBs) 
sowie die PowerPoint-Präsentationen für die Hörsäle und konvertiert diese direkt in PDFs.
""")

st.divider()

# --- KONFIGURATION & USER-PFADE ---
st.subheader("⚙️ 1. Basiseinstellungen")
col_config1, col_config2 = st.columns(2)

with col_config1:
    semester_input = st.text_input("Aktuelles Semester (für die Folien):", value="Sommersemester 2026")

with col_config2:
    # Standardmäßig wird der Ordner dort erstellt, wo das Skript läuft
    default_output = os.path.join(os.getcwd(), "lve_online")
    output_root_input = st.text_input("Ausgabe-Verzeichnis (Wo soll die Struktur entstehen?):", value=default_output)

st.divider()

# --- DYNAMISCHER DATEN-UPLOAD BEREICH ---
st.subheader("📂 2. Benötigte Dokumente & Vorlagen hochladen")
st.info("Lade hier alle Dateien flexibel hoch – unabhängig davon, wo sie auf deinem PC gespeichert sind.")

up_col1, up_col2 = st.columns(2)

with up_col1:
    st.markdown("### 📄 Basisdaten & Struktur")
    uploaded_json = st.file_uploader("Strukturdatei hochladen (z.B. kombiniert_mit_losungen.json)", type=["json"])
    uploaded_logo = st.file_uploader("Universitäts-Logo hochladen (PNG)", type=["png"])
    
    st.markdown("### 🌐 HTML-Vorlagen (RLB)")
    uploaded_html_def = st.file_uploader("Standard-HTML Vorlage (z.B. Leere_RLB.html)", type=["html"])
    uploaded_html_fsr = st.file_uploader("FSR-HTML Vorlage (z.B. Leere_FSR_RLB.html)", type=["html"])

with up_col2:
    st.markdown("### 🗂️ Fakultäts-Infoblätter (PDF)")
    uploaded_info_mathe = st.file_uploader("Infoblatt für Mathematik (PDF)", type=["pdf"])
    uploaded_info_wiwi = st.file_uploader("Infoblatt für Wirtschaftswissenschaften (PDF)", type=["pdf"])
    uploaded_info_allg = st.file_uploader("Allgemeines Infoblatt / restliche Teilbereiche (PDF)", type=["pdf"])

st.divider()

# --- MAILING & KENNUNGSMAPPINGS ---
TEILBEREICHE = [
    "Biologie", "Chemie", "Humboldt Studienzentrum", "Informatik", 
    "Ingenieurwissenschaften", "Mathematik", "Physik", "Psychologie", 
    "Wirtschaftswissenschaften", "Zentrum für Sprachen und Philologie"
]

ART_MAPPING = {
    "Vorlesung": "1", "Seminar": "2", "Übung": "4", "Praktikum": "5",
    "Vorlesung/Seminar": "8", "Vorlesung/Übung": "12", "Tutorium": "31",
    "Online-Seminar": "41", "Vorlesung/Übung/Tutorium": "43", "Klinisches Praktikum": "44"
}

# --- HILFSFUNKTIONEN FÜR DIE DYNAMISCHE GENERIERUNG ---
def clean_filename(name):
    if not name: return "Unbekannt"
    name = str(name).replace('ä', 'ae').replace('ö', 'oe').replace('ü', 'ue').replace('ß', 'ss')
    cleaned = re.sub(r'[\\/:*?"<>|&#~,;+"!§$^°\[\]{}=]', "_", name)
    return re.sub(r'_+', "_", cleaned).strip("_ ")

def get_last_name(full_name):
    parts = str(full_name).split()
    return parts[-1] if parts else "Unbekannt"

def should_add_suffix(titel, art_str):
    if not art_str: return False
    return art_str.lower() not in titel.lower()

def process_html_template(html_content_bytes, output_path, data, v_art_str):
    try:
        content = html_content_bytes.decode("utf-8")
        titel_raw = data.get("lve_titel", "")
        titel_anzeige = f"{titel_raw} ({v_art_str})" if should_add_suffix(titel_raw, v_art_str) else titel_raw
        
        replacements = [
            ("semester", semester_input),
            ("titel", titel_anzeige),
            ("kennung", data.get("lve_kennung", "")),
            ("dozent", data.get("prof_name", "")),
            ("dateiname_RLB", os.path.basename(output_path))
        ]
        for field_id, value in replacements:
            pattern = rf'(<[^>]*id="{field_id}"[^>]*value=")([^"]*)(")'
            content = re.sub(pattern, rf'\1{value}\3', content)
            
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)
    except Exception as e:
        raise Exception(f"HTML Verarbeitungsfehler: {e}")

def create_complex_ppt(output_path, qr_code_path, logo_path, data, prof_name, suffix_clean):
    veranstaltung_pure = data.get("lve_titel", "Unbekannt")
    veranstaltung_mit_suffix = f"{veranstaltung_pure} ({suffix_clean})" if suffix_clean and suffix_clean.lower() not in veranstaltung_pure.lower() else veranstaltung_pure
    login_losung = data.get("losung", "")
    qr_code_url = f"https://evaluation.uni-ulm.de/evasys/online.php?p={login_losung}"

    prs = Presentation()
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    text_box = slide.shapes.add_textbox(Cm(0.1), Cm(0), Cm(8), Cm(1))
    text_box.text_frame.text = semester_input
    text_box.text_frame.paragraphs[0].runs[0].font.bold = True
    
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Cm(18), Cm(0.2), width=Cm(7))

    shape = slide.shapes.add_table(4, 2, Cm(1.25), Cm(3), Cm(23), Cm(5))
    table = shape.table
    table.columns[0].width, table.columns[1].width = Cm(7), Cm(16)
    
    try:
        shape._element.graphic.graphicData.tbl[0][-1].text = '{7E9639D4-E3E2-4D34-9284-5A2195B3D0D7}'
    except: pass
    
    headers = ["Lehrveranstaltung:", "Lehrperson:", "URL:", "TAN/Losung:"]
    contents = [veranstaltung_mit_suffix, prof_name, "https://evaluation.uni-ulm.de/evasys/online", login_losung]

    for row in range(4):
        cell_label = table.cell(row, 0)
        cell_label.fill.solid()
        cell_label.fill.fore_color.rgb = RGBColor(125, 154, 170)
        cell_label.text = headers[row]
        cell_label.text_frame.paragraphs[0].runs[0].font.bold = True
        cell_label.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell_label.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        cell_val = table.cell(row, 1)
        p = cell_val.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        cell_val.vertical_anchor = MSO_ANCHOR.MIDDLE
        r = p.add_run()
        r.text = contents[row]
        if row >= 2: r.hyperlink.address = qr_code_url

    slide.shapes.add_textbox(Cm(3.5), Cm(14), Cm(8), Cm(1)).text_frame.text = "QR-Code zum Scannen:"
    slide.shapes.add_picture(qr_code_path, Cm(13.32), Cm(11.5), width=Cm(6), height=Cm(6))
    
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.add_textbox(Cm(0.1), Cm(0), Cm(8), Cm(1)).text_frame.text = semester_input
    slide2.shapes.add_textbox(Cm(0.1), Cm(0.8), Cm(17), Cm(6.8)).text_frame.text = f"{veranstaltung_mit_suffix}\nLehrperson: {prof_name}"
    
    if logo_path and os.path.exists(logo_path):
        slide2.shapes.add_picture(logo_path, Cm(18), Cm(0.2), width=Cm(7))
    slide2.shapes.add_picture(qr_code_path, Cm(6.3), Cm(5.6), width=Cm(12.78), height=Cm(12.78))

    prs.save(output_path)

def convert_to_pdf(target_root):
    try:
        import comtypes.client
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1
        for root, _, files in os.walk(target_root):
            for f in files:
                if f.endswith(".pptx") and not f.startswith("~$"):
                    path = os.path.normpath(os.path.abspath(os.path.join(root, f)))
                    pdf_path = os.path.splitext(path)[0] + ".pdf"
                    if not os.path.exists(pdf_path):
                        pres = powerpoint.Presentations.Open(path, WithWindow=0, ReadOnly=1)
                        pres.ExportAsFixedFormat(pdf_path, 2, 1, 0, 1, 1, 0, None, 1, "", True, True, True, False)
                        pres.Close()
        powerpoint.Quit()
    except Exception as e:
        st.warning(f"⚠️ PDF-Konvertierung über PowerPoint fehlgeschlagen: {e}. Die PowerPoint-Dateien (.pptx) wurden trotzdem generiert. Hinweis: Auf Cloud-Servern ohne Windows/Office schlägt dieser Teilschritt systembedingt fehl.")

# --- EXECUTION BUTTON ---
st.subheader("🚀 3. Prozess starten")

all_files_uploaded = all([
    uploaded_json, uploaded_logo, uploaded_html_def, uploaded_html_fsr,
    uploaded_info_mathe, uploaded_info_wiwi, uploaded_info_allg
])

if not all_files_uploaded:
    st.error("❌ Bitte lade zuerst alle oben geforderten 7 Dateien hoch.")
else:
    if st.button("Unterlagen jetzt vollautomatisch generieren", type="primary"):
        progress_bar = st.progress(0)
        status_text = st.empty()
        
        try:
            # Temporären Cache-Ordner für Medien anlegen
            tmp_dir = os.path.join(os.getcwd(), "tmp_generator_cache")
            os.makedirs(tmp_dir, exist_ok=True)
            
            # Speicher Medien temporär ab
            tmp_logo_path = os.path.join(tmp_dir, "logo.png")
            tmp_pdf_mathe = os.path.join(tmp_dir, "mathe.pdf")
            tmp_pdf_wiwi = os.path.join(tmp_dir, "wiwi.pdf")
            tmp_pdf_allg = os.path.join(tmp_dir, "allgemein.pdf")
            
            with open(tmp_logo_path, "wb") as f: f.write(uploaded_logo.getbuffer())
            with open(tmp_pdf_mathe, "wb") as f: f.write(uploaded_info_mathe.getbuffer())
            with open(tmp_pdf_wiwi, "wb") as f: f.write(uploaded_info_wiwi.getbuffer())
            with open(tmp_pdf_allg, "wb") as f: f.write(uploaded_info_allg.getbuffer())
            
            # Ausgabe-Verzeichnis leeren/erstellen
            os.makedirs(output_root_input, exist_ok=True)

            # JSON Daten parsen
            data_struct = json.loads(uploaded_json.read().decode("utf-8"))
            
            total_tb = len([x for x in TEILBEREICHE if x in data_struct])
            current_tb = 0

            # Hauptschleife über alle Fachbereiche
            for tb in TEILBEREICHE:
                if tb not in data_struct: continue
                current_tb += 1
                status_text.text(f"⏳ Verarbeite Fachbereich: {tb} ({current_tb}/{total_tb})...")
                
                for prof in data_struct[tb]:
                    prof_full = prof["prof_name"]
                    nachname_clean = clean_filename(get_last_name(prof_full))
                    prof_folder = os.path.join(output_root_input, tb, nachname_clean)
                    os.makedirs(prof_folder, exist_ok=True)
                    
                    for match in prof["matches"]:
                        if not match.get("losung"): continue

                        kennung_clean = clean_filename(match["lve_kennung"])
                        titel_clean = clean_filename(match["lve_titel"])
                        
                        v_art_id = str(match["codes"]["veranstaltungsart_id"])
                        v_art_str = next((k for k, v in ART_MAPPING.items() if v == v_art_id), "")
                        v_art_clean = clean_filename(v_art_str)
                        
                        v_art_suffix = f"_{v_art_clean}" if v_art_clean and not titel_clean.endswith(v_art_clean) else ""
                        titel_gekürzt = titel_clean[:80].strip("_ ") if len(titel_clean) > 80 else titel_clean
                        if len(titel_clean) > 80 and v_art_clean: v_art_suffix = f"_{v_art_clean}"

                        file_base = f"{nachname_clean}_{kennung_clean}_{titel_gekürzt}{v_art_suffix}"
                        
                        # --- HTML / RLB ERSTELLUNG ---
                        html_path = os.path.join(prof_folder, f"{nachname_clean}_RLB_{kennung_clean}_{titel_gekürzt}{v_art_suffix}.html")
                        raw_html_bytes = uploaded_html_fsr.getvalue() if tb in ["Mathematik", "Wirtschaftswissenschaften"] else uploaded_html_def.getvalue()
                        process_html_template(raw_html_bytes, html_path, match, v_art_str)
                        
                        # --- QR-CODE GENERIEREN ---
                        qr_temp = os.path.join(prof_folder, f"temp_{kennung_clean}.png")
                        qr = qrcode.QRCode(version=None, border=1)
                        qr.add_data(f"https://evaluation.uni-ulm.de/evasys/online.php?p={match['losung']}")
                        qr.make(fit=True)
                        qr.make_image().save(qr_temp)
                        
                        # --- POWERPOINT ERSTELLUNG ---
                        ppt_path = os.path.join(prof_folder, f"{file_base}.pptx")
                        create_complex_ppt(ppt_path, qr_temp, tmp_logo_path, match, prof_full, v_art_str)
                        if os.path.exists(qr_temp): os.remove(qr_temp)

                progress_bar.progress(int((current_tb / total_tb) * 75))

            # --- PDF-KONVERTIERUNG ---
            status_text.text("⏳ Konvertiere PowerPoint-Folien in PDFs...")
            convert_to_pdf(output_root_input)
            progress_bar.progress(90)
            
            # --- DYNAMISCHES INFOPLATT-KOPIEREN ---
            status_text.text("⏳ Verteile hochgeladene Infoblätter in Dozentenordner...")
            for tb_name in os.listdir(output_root_input):
                tb_path = os.path.join(output_root_input, tb_name)
                if not os.path.isdir(tb_path): continue
                
                if tb_name == "Mathematik": src_pdf = tmp_pdf_mathe
                elif tb_name == "Wirtschaftswissenschaften": src_pdf = tmp_pdf_wiwi
                else: src_pdf = tmp_pdf_allg
                
                orig_name = uploaded_info_mathe.name if tb_name == "Mathematik" else (uploaded_info_wiwi.name if tb_name == "Wirtschaftswissenschaften" else uploaded_info_allg.name)

                for prof in os.listdir(tb_path):
                    prof_dir = os.path.join(tb_path, prof)
                    if os.path.isdir(prof_dir):
                        shutil.copy2(src_pdf, os.path.join(prof_dir, orig_name))
            
            # Cache aufräumen
            shutil.rmtree(tmp_dir, ignore_errors=True)
            
            progress_bar.progress(100)
            status_text.empty()
            st.success(f"🏁 Ausführung erfolgreich beendet! Alle Unterlagen liegen bereit unter: `{output_root_input}`")
            
        except Exception as e:
            st.error(f"❌ Laufzeitfehler während der Generierung: {e}")
